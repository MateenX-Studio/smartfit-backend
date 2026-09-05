
# Import FastAPI to create our web server
from fastapi import FastAPI, UploadFile, File, Form, Header, HTTPException, Depends
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import io
import os
import json
from dotenv import load_dotenv
from google import genai
from supabase import create_client
from pydantic import BaseModel
from fastapi.middleware.cors import CORSMiddleware
import time


# Load environment variables (API keys)
load_dotenv()

# Set up Gemini client
api_key = os.getenv("GEMINI_API_KEY")
client = genai.Client(api_key=api_key)

# Set up Supabase client
supabase_url = os.getenv("SUPABASE_URL")
supabase_key = os.getenv("SUPABASE_KEY")
supabase = create_client(supabase_url, supabase_key)

# Create the FastAPI application instance
app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    # allow_origins=["http://localhost:3000"],
    allow_origins=["*"],  # Allow all origins (Vercel, localhost, etc.)
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Verify the access token and return the user's ID
def get_current_user(authorization: str = Header(...)):
    try:
        token = authorization.replace("Bearer ", "")
        user_response = supabase.auth.get_user(token)
        return user_response.user.id
    except Exception:
        raise HTTPException(status_code=401, detail="Invalid or expired token")


# ---------- File text extraction helpers ----------

def extract_text_from_pdf(file_bytes):
    pdf_reader = PdfReader(io.BytesIO(file_bytes))
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text


def extract_text_from_docx(file_bytes):
    doc = Document(io.BytesIO(file_bytes))
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text += cell.text + "\n"
    return text


def extract_text_from_pptx(file_bytes):
    presentation = Presentation(io.BytesIO(file_bytes))
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text_frame.text + "\n"
    return text


# ---------- Auth models ----------

class AuthRequest(BaseModel):
    email: str
    password: str


# ---------- Routes ----------

@app.get("/")
def read_root():
    return {"message": "SmartFit AI backend is running!"}

# @app.post("/signup")
# def signup(auth: AuthRequest):
#     try:
#         response = supabase.auth.sign_up({
#             "email": auth.email,
#             "password": auth.password
#         })
#         return {"message": "Signup successful", "user_id": response.user.id}
#     except Exception as e:
#         raise HTTPException(status_code=400, detail=str(e))
@app.post("/signup")
def signup(auth: AuthRequest):
    try:
        response = supabase.auth.sign_up({
            "email": auth.email,
            "password": auth.password
        })
        return {
            "message": "Signup successful",
            "user_id": response.user.id,
            "access_token": response.session.access_token if response.session else None,
            "refresh_token": response.session.refresh_token if response.session else None,
        }
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/login")
def login(auth: AuthRequest):
    response = supabase.auth.sign_in_with_password({
        "email": auth.email,
        "password": auth.password
    })
    return {
        "message": "Login successful",
        "access_token": response.session.access_token,
        "refresh_token": response.session.refresh_token,
        "user_id": response.user.id
    }




@app.post("/analyze")
async def analyze_cv(
    cv_file: UploadFile = File(...),
    job_description: str = Form(...),
    user_id: str = Depends(get_current_user)
):
    file_bytes = await cv_file.read()
    filename = cv_file.filename.lower()

    if filename.endswith(".pdf"):
        cv_text = extract_text_from_pdf(file_bytes)
    elif filename.endswith(".docx"):
        cv_text = extract_text_from_docx(file_bytes)
    elif filename.endswith(".pptx"):
        cv_text = extract_text_from_pptx(file_bytes)
    else:
        return {"error": "Unsupported file type. Please upload a PDF, DOCX, or PPTX file."}

    prompt = f"""You are a professional career advisor. Analyze the following CV against the job description.

CV:
{cv_text}

Job Description:
{job_description}

Respond with ONLY a valid JSON object (no extra text, no markdown formatting) in this exact structure:
{{
  "match_score": <number from 0 to 100>,
  "missing_skills": ["skill1", "skill2"],
  "suggestions": ["suggestion1", "suggestion2", "suggestion3"],
  "suggested_roles": ["role1", "role2", "role3"]
}}
"""
    # Try calling Gemini up to 3 times if the server is temporarily busy
    max_retries = 3
    response = None
    last_error = None

    for attempt in range(max_retries):
        try:
            response = client.models.generate_content(
                model="gemini-3.6-flash",
                contents=prompt
            )
            break  # success, stop retrying
        except Exception as e:
            last_error = e
            if attempt < max_retries - 1:
                time.sleep(2)  # wait 2 seconds before trying again

    if response is None:
        raise HTTPException(
            status_code=503,
            detail="The AI service is temporarily busy. Please try again in a moment."
        )
  

    raw_text = response.text.strip()
    if raw_text.startswith("```"):
        raw_text = raw_text.replace("```json", "").replace("```", "").strip()

    try:
        result = json.loads(raw_text)
    except json.JSONDecodeError:
        return {"error": "Failed to parse AI response", "raw": raw_text}

    supabase.table("analyses").insert({
        "user_id": user_id,
        "job_description": job_description,
        "match_score": result.get("match_score"),
        "missing_skills": ", ".join(result.get("missing_skills", [])),
        "suggestions": ", ".join(result.get("suggestions", [])),
        "suggested_roles": ", ".join(result.get("suggested_roles", []))
    }).execute()

    return result    

@app.get("/history")
def get_history(user_id: str = Depends(get_current_user)):
    response = supabase.table("analyses").select("*").eq("user_id", user_id).order("created_at", desc=True).execute()
    return response.data    
class RefreshRequest(BaseModel):
    refresh_token: str

@app.post("/refresh")
def refresh_token(req: RefreshRequest):
    try:
        response = supabase.auth.refresh_session(req.refresh_token)
        return {
            "access_token": response.session.access_token,
            "refresh_token": response.session.refresh_token,
        }
    except Exception as e:
        raise HTTPException(status_code=401, detail="Session expired. Please log in again.")